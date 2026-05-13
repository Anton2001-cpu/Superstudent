import json
import os
import hashlib
import re
from pathlib import Path

import numpy as np
from openai import OpenAI
import pypdf
from docx import Document as DocxDocument


# ── Supabase vector store ────────────────────────────────────────────────────

class _SupabaseCollection:
    def __init__(self, client):
        self._sb = client

    def count(self):
        result = self._sb.table("chunks").select("id", count="exact").execute()
        return result.count or 0

    def upsert(self, ids, embeddings, documents, metadatas):
        rows = [
            {
                "id": id_,
                "embedding": emb,
                "document": doc,
                "course": meta["course"],
                "filename": meta["filename"],
                "location": meta.get("location", ""),
                "preview": meta.get("preview", ""),
            }
            for id_, emb, doc, meta in zip(ids, embeddings, documents, metadatas)
        ]
        for i in range(0, len(rows), 100):
            self._sb.table("chunks").upsert(rows[i:i + 100]).execute()

    def query(self, query_embeddings, n_results=3, where=None, include=None):
        filter_course = None
        filter_filenames = None
        if where:
            if "$or" in where:
                filter_filenames = [c["filename"]["$eq"] for c in where["$or"]]
            elif "filename" in where:
                filter_filenames = [where["filename"]["$eq"]]
            elif "course" in where:
                filter_course = where["course"]["$eq"]

        result = self._sb.rpc("match_chunks", {
            "query_embedding": query_embeddings[0],
            "match_count": n_results,
            "filter_course": filter_course,
            "filter_filenames": filter_filenames,
        }).execute()

        rows = result.data or []
        return {
            "ids": [[r["id"] for r in rows]],
            "documents": [[r["document"] for r in rows]],
            "metadatas": [[{
                "course": r["course"], "filename": r["filename"],
                "location": r["location"], "preview": r["preview"],
            } for r in rows]],
            "distances": [[1 - r["similarity"] for r in rows]],
        }

    def get(self, where=None, include=None):
        q = self._sb.table("chunks").select("id,document,course,filename,location,preview")
        if where:
            if "$and" in where:
                for cond in where["$and"]:
                    for key, val in cond.items():
                        q = q.eq(key, val["$eq"])
            elif "course" in where:
                q = q.eq("course", where["course"]["$eq"])
        result = q.limit(500).execute()
        rows = result.data or []
        return {
            "ids": [r["id"] for r in rows],
            "documents": [r["document"] for r in rows],
            "metadatas": [{
                "course": r["course"], "filename": r["filename"],
                "location": r["location"], "preview": r["preview"],
            } for r in rows],
        }

    def delete(self, ids):
        for i in range(0, len(ids), 100):
            self._sb.table("chunks").delete().in_("id", ids[i:i + 100]).execute()


# ── In-memory fallback (local dev without Supabase) ─────────────────────────

class _LocalCollection:
    def __init__(self):
        self._ids = []
        self._embeddings = []
        self._documents = []
        self._metadatas = []

    def count(self):
        return len(self._ids)

    def upsert(self, ids, embeddings, documents, metadatas):
        existing = {id_: i for i, id_ in enumerate(self._ids)}
        for id_, emb, doc, meta in zip(ids, embeddings, documents, metadatas):
            if id_ in existing:
                i = existing[id_]
                self._embeddings[i] = emb
                self._documents[i] = doc
                self._metadatas[i] = meta
            else:
                self._ids.append(id_)
                self._embeddings.append(emb)
                self._documents.append(doc)
                self._metadatas.append(meta)

    def query(self, query_embeddings, n_results=3, where=None, include=None):
        if not self._ids:
            return {"ids": [[]], "documents": [[]], "metadatas": [[]], "distances": [[]]}
        q = np.array(query_embeddings[0], dtype=np.float32)
        embs = np.array(self._embeddings, dtype=np.float32)
        q_norm = q / (np.linalg.norm(q) + 1e-10)
        embs_norm = embs / (np.linalg.norm(embs, axis=1, keepdims=True) + 1e-10)
        similarities = embs_norm @ q_norm
        indices = list(range(len(self._ids)))
        if where:
            indices = [i for i in indices if self._match(self._metadatas[i], where)]
        indices.sort(key=lambda i: float(similarities[i]), reverse=True)
        top = indices[:n_results]
        return {
            "ids": [[self._ids[i] for i in top]],
            "documents": [[self._documents[i] for i in top]],
            "metadatas": [[self._metadatas[i] for i in top]],
            "distances": [[float(1 - similarities[i]) for i in top]],
        }

    def get(self, where=None, include=None):
        if where is None:
            return {"ids": list(self._ids), "documents": list(self._documents), "metadatas": list(self._metadatas)}
        indices = [i for i in range(len(self._ids)) if self._match(self._metadatas[i], where)]
        return {
            "ids": [self._ids[i] for i in indices],
            "documents": [self._documents[i] for i in indices],
            "metadatas": [self._metadatas[i] for i in indices],
        }

    def delete(self, ids):
        id_set = set(ids)
        keep = [i for i in range(len(self._ids)) if self._ids[i] not in id_set]
        self._ids = [self._ids[i] for i in keep]
        self._embeddings = [self._embeddings[i] for i in keep]
        self._documents = [self._documents[i] for i in keep]
        self._metadatas = [self._metadatas[i] for i in keep]

    def _match(self, meta, where):
        if "$and" in where:
            return all(self._match(meta, c) for c in where["$and"])
        if "$or" in where:
            return any(self._match(meta, c) for c in where["$or"])
        for key, condition in where.items():
            val = meta.get(key)
            if isinstance(condition, dict):
                if "$eq" in condition and val != condition["$eq"]:
                    return False
            elif val != condition:
                return False
        return True


# ── RAG Engine ────────────────────────────────────────────────────────────────

class RAGEngine:
    def __init__(self):
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.chunk_size = 800
        self.chunk_overlap = 100

        supabase_url = os.getenv("SUPABASE_URL")
        supabase_key = os.getenv("SUPABASE_KEY")
        if supabase_url and supabase_key:
            from supabase import create_client
            self._sb = create_client(supabase_url, supabase_key)
            self.collection = _SupabaseCollection(self._sb)
        else:
            self._sb = None
            self.collection = _LocalCollection()
            self._load_seed_data()

    def _load_seed_data(self):
        seed_file = Path(__file__).parent / "seed_data.json"
        if not seed_file.exists():
            return
        with open(seed_file, "r", encoding="utf-8") as f:
            data = json.load(f)
        self.collection.upsert(
            ids=data["ids"],
            embeddings=data["embeddings"],
            documents=data["documents"],
            metadatas=data["metadatas"],
        )

    # ── Text extraction ───────────────────────────────────────────────────

    def extract_with_meta(self, file_path: str) -> list:
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return self._extract_docx(file_path)
        elif ext == ".txt":
            return self._extract_txt(file_path)
        elif ext in (".pptx", ".ppsx"):
            return self._extract_pptx(file_path)
        elif ext in (".mp4", ".mp3", ".m4a"):
            return self._extract_mp4(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _extract_mp4(self, file_path: str) -> list:
        import io
        with open(file_path, "rb") as f:
            audio_bytes = f.read()
        if len(audio_bytes) > 25 * 1024 * 1024:
            raise ValueError("Videobestand is te groot voor transcriptie (max 25 MB). Knip het bestand eerst bij.")
        bio = io.BytesIO(audio_bytes)
        bio.name = Path(file_path).name
        transcript = self.client.audio.transcriptions.create(
            model="whisper-1",
            file=bio,
            response_format="text",
        )
        text = transcript.strip() if isinstance(transcript, str) else str(transcript).strip()
        if not text:
            raise ValueError("Geen spraak gevonden in het videobestand.")
        return [{"text": text, "location": "Transcriptie"}]

    def _is_title_page(self, text: str, page_index: int) -> bool:
        if page_index != 0:
            return False
        if len(text.strip()) < 600:
            return True
        t = text.lower()
        has_email = bool(re.search(r'[\w.+-]+@[\w-]+\.[a-z]{2,}', text))
        has_abstract = "abstract" in t[:800]
        has_affiliation = any(w in t[:600] for w in [
            "university", "universiteit", "institute", "department", "faculty",
            "corresponding author", "delft", "technische"
        ])
        if has_email and (has_abstract or has_affiliation):
            return True
        return False

    def _is_reference_page(self, text: str) -> bool:
        first_lines = text.strip()[:300].lower()
        ref_headings = ["references", "bibliography", "referenties", "literatuur",
                        "bronnen", "literatuurlijst", "works cited", "bibliographie"]
        if any(first_lines.startswith(h) or f"\n{h}" in first_lines for h in ref_headings):
            return True
        citation_hits = len(re.findall(r"^\s*\[\d+\]", text, re.MULTILINE))
        if citation_hits >= 4:
            return True
        return False

    def _extract_pdf(self, path: str) -> list:
        reader = pypdf.PdfReader(path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if not text or not text.strip():
                continue
            if self._is_title_page(text, i):
                continue
            if self._is_reference_page(text):
                continue
            pages.append({"text": text.strip(), "location": f"Page {i + 1}"})
        return pages

    def _extract_docx(self, path: str) -> list:
        doc = DocxDocument(path)
        sections = []
        current_lines = []
        section_num = 1
        current_heading = f"Section {section_num}"

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith("Heading"):
                if current_lines:
                    sections.append({"text": "\n".join(current_lines), "location": current_heading})
                    current_lines = []
                    section_num += 1
                current_heading = text or f"Section {section_num}"
                current_lines.append(text)
            else:
                current_lines.append(text)

        if current_lines:
            sections.append({"text": "\n".join(current_lines), "location": current_heading})

        if not sections:
            full = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return [{"text": full, "location": "Section 1"}] if full else []

        return sections

    def _extract_pptx(self, path: str) -> list:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            text = "\n".join(texts).strip()
            if text:
                slides.append({"text": text, "location": f"Slide {i + 1}"})
        return slides

    def _extract_txt(self, path: str) -> list:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        lines = text.split("\n")
        blocks = []
        size = 60
        for i in range(0, len(lines), size):
            chunk = "\n".join(lines[i: i + size]).strip()
            if chunk:
                end = min(i + size, len(lines))
                blocks.append({"text": chunk, "location": f"Lines {i + 1}-{end}"})
        return blocks or [{"text": text.strip(), "location": "Full text"}]

    # ── Chunking ──────────────────────────────────────────────────────────

    def chunk_section(self, text: str, location: str) -> list:
        chunks = []
        start = 0
        text = text.strip()
        part = 1
        while start < len(text):
            end = start + self.chunk_size
            chunk = text[start:end].strip()
            if chunk:
                loc = location if len(text) <= self.chunk_size else f"{location}, part {part}"
                chunks.append({"text": chunk, "location": loc})
                part += 1
            start += self.chunk_size - self.chunk_overlap
        return chunks

    # ── Document ingestion ────────────────────────────────────────────────

    def add_document(self, file_path: str, course_name: str, filename: str) -> int:
        sections = self.extract_with_meta(file_path)
        if not sections:
            raise ValueError("Could not extract any text from this file.")

        all_chunks = []
        for section in sections:
            all_chunks.extend(self.chunk_section(section["text"], section["location"]))

        if not all_chunks:
            raise ValueError("File appears to be empty after processing.")

        texts = [c["text"] for c in all_chunks]

        all_embeddings = []
        for i in range(0, len(texts), 100):
            batch = texts[i: i + 100]
            response = self.client.embeddings.create(input=batch, model="text-embedding-3-small")
            all_embeddings.extend([item.embedding for item in response.data])

        ids = [
            hashlib.sha256(f"{course_name}|{filename}|{i}".encode()).hexdigest()
            for i in range(len(all_chunks))
        ]

        metadatas = [
            {
                "course": course_name,
                "filename": filename,
                "location": c["location"],
                "preview": c["text"][:900],
            }
            for c in all_chunks
        ]

        self.collection.upsert(ids=ids, embeddings=all_embeddings, documents=texts, metadatas=metadatas)

        ext = Path(file_path).suffix.lower()
        if self._sb:
            if ext in (".pptx", ".ppsx"):
                self._store_pptx_images(file_path, course_name, filename)
            elif ext == ".pdf":
                self._store_pdf_images(file_path, course_name, filename)

        return len(all_chunks)

    def _store_pptx_images(self, file_path: str, course_name: str, filename: str):
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return
        try:
            prs = Presentation(file_path)
        except Exception:
            return
        try:
            self._sb.table("slide_images").delete().eq("course", course_name).eq("filename", filename).execute()
        except Exception:
            pass
        rows = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            img_idx = 0
            for shape in slide.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        if len(image.blob) < 5000:
                            continue
                        ext = (image.ext or "png").lstrip(".")
                        storage_path = f"{course_name}/{filename}/slide_{slide_num}_img_{img_idx}.{ext}"
                        self._sb.storage.from_("course-files").upload(
                            storage_path, image.blob,
                            {"content-type": image.content_type or "image/png", "upsert": "true"},
                        )
                        rows.append({
                            "course": course_name,
                            "filename": filename,
                            "slide": f"Slide {slide_num}",
                            "storage_path": storage_path,
                            "image_index": img_idx,
                        })
                        img_idx += 1
                except Exception:
                    pass
        if rows:
            try:
                self._sb.table("slide_images").insert(rows).execute()
            except Exception:
                pass

    def _store_pdf_images(self, file_path: str, course_name: str, filename: str):
        try:
            reader = pypdf.PdfReader(file_path)
        except Exception:
            return
        try:
            self._sb.table("slide_images").delete().eq("course", course_name).eq("filename", filename).execute()
        except Exception:
            pass
        rows = []
        seen_hashes = set()
        for page_num, page in enumerate(reader.pages, start=1):
            img_idx = 0
            try:
                page_images = list(page.images)
            except Exception:
                continue
            for image in page_images:
                if img_idx >= 3:
                    break
                try:
                    data = image.data
                    if not data or len(data) < 5000:
                        continue
                    img_hash = hashlib.md5(data).hexdigest()
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)
                    name = image.name or ""
                    ext = name.rsplit(".", 1)[-1].lower() if "." in name else "png"
                    if ext not in ("jpg", "jpeg", "png", "gif", "webp"):
                        ext = "png"
                    mime = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                            "gif": "image/gif", "webp": "image/webp"}.get(ext, "image/png")
                    storage_path = f"{course_name}/{filename}/page_{page_num}_img_{img_idx}.{ext}"
                    self._sb.storage.from_("course-files").upload(
                        storage_path, data,
                        {"content-type": mime, "upsert": "true"},
                    )
                    rows.append({
                        "course": course_name,
                        "filename": filename,
                        "slide": f"Page {page_num}",
                        "storage_path": storage_path,
                        "image_index": img_idx,
                    })
                    img_idx += 1
                except Exception:
                    pass
        if rows:
            try:
                self._sb.table("slide_images").insert(rows).execute()
            except Exception:
                pass

    def get_slide_images(self, metas: list) -> list:
        if not self._sb:
            return []
        source_locs = set()
        for meta in metas:
            fname = meta.get("filename", "").lower()
            if fname.endswith(".pptx") or fname.endswith(".ppsx") or fname.endswith(".pdf"):
                location = meta.get("location", "")
                base_loc = location.split(",")[0].strip()
                source_locs.add((meta.get("course", ""), meta.get("filename", ""), base_loc))
        if not source_locs:
            return []
        urls = []
        for course, filename, slide in source_locs:
            try:
                rows = (
                    self._sb.table("slide_images")
                    .select("storage_path")
                    .eq("course", course)
                    .eq("filename", filename)
                    .eq("slide", slide)
                    .execute()
                    .data or []
                )
                for row in rows:
                    try:
                        signed = self._sb.storage.from_("course-files").create_signed_url(row["storage_path"], 3600)
                        if isinstance(signed, dict):
                            url = (signed.get("signedURL") or signed.get("signedUrl")
                                   or (signed.get("data") or {}).get("signedUrl"))
                        else:
                            url = getattr(signed, "signed_url", None)
                        if url:
                            urls.append(url)
                    except Exception:
                        pass
            except Exception:
                pass
        return urls

    # ── Query (student) ───────────────────────────────────────────────────

    def _expand_query(self, text: str) -> str:
        variants = {text}
        variants.add(re.sub(r'(\w)-(\w)', lambda m: f"{m.group(1)} {m.group(2)}", text))
        variants.add(re.sub(r'(\b\w+)\s+(\w+\b)', lambda m: f"{m.group(1)}-{m.group(2)}", text))
        return " | ".join(v for v in variants if v != text) or text

    def query(self, question: str, books: list = None, history: list = None, lang: str = "EN") -> dict:
        total = self.collection.count()
        if total == 0:
            return {
                "answer": "I couldn't find anything — no course materials have been uploaded yet. Please ask your teacher.",
                "sources": [],
            }

        search_query = self._get_search_query(question, history or [])
        response = self.client.embeddings.create(
            input=[self._expand_query(search_query)],
            model="text-embedding-3-small",
        )
        query_embedding = response.data[0].embedding

        where = None
        if books and len(books) > 0:
            if len(books) == 1:
                where = {"filename": {"$eq": books[0]}}
            else:
                where = {"$or": [{"filename": {"$eq": b}} for b in books]}

        n_results = min(2, total)

        results = self.collection.query(query_embeddings=[query_embedding], n_results=n_results, where=where)

        docs = results["documents"][0]
        metas = results["metadatas"][0]

        if not docs:
            return {
                "answer": "I couldn't find this in your course material, click below to check relevant information online.",
                "sources": [],
            }

        context_parts = []
        sources = []
        for doc, meta in zip(docs, metas):
            fname = meta.get("filename", "Unknown file")
            loc = meta.get("location", "")
            context_parts.append(f"[{fname} — {loc}]\n{doc[:700]}")
            sources.append({
                "filename": fname,
                "location": loc,
                "preview": meta.get("preview", doc[:500]),
            })

        context = "\n\n---\n\n".join(context_parts)

        lang_instruction = "Answer in Dutch." if lang == "NL" else "Answer in English."
        fallback_phrase = (
            "Ik kon dit niet vinden in je cursusmateriaal, klik hieronder voor relevante informatie online."
            if lang == "NL" else
            "I couldn't find this in your course material, click below to check relevant information online."
        )
        answer_prompt = (
            f"You are a study assistant. {lang_instruction} "
            "Answer using ONLY the course material below. No emojis, no links, no outside knowledge. "
            "Use bullet points for lists. Bold key terms with **term**. 4-8 sentences. "
            f"If not found, say exactly: \"{fallback_phrase}\""
        )

        messages = [{"role": "system", "content": answer_prompt}]
        if history:
            messages.extend(history[-4:])
        messages.append({"role": "user", "content": f"Material:\n{context}\n\nQuestion: {question}"})

        chat_response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0,
            max_tokens=400,
        )
        answer = chat_response.choices[0].message.content.strip()

        fallback_old = "I couldn't find this in your course materials. Please ask your teacher."
        if fallback_old in answer and len(answer) > len(fallback_old) + 10:
            answer = answer.replace(fallback_old, "").strip().rstrip("\n").strip()
        no_answer = "couldn't find" in answer.lower() or "kon dit niet vinden" in answer.lower()

        extra = ""
        if no_answer:
            reply_lang = "Dutch" if lang == "NL" else "English"
            try:
                extra_response = self.client.responses.create(
                    model="gpt-4o-mini",
                    tools=[{"type": "web_search_preview"}],
                    input=(
                        f"The student asked: \"{question}\". "
                        "Search online and give a short, helpful answer (2-4 sentences) about this topic. "
                        f"You MUST respond entirely in {reply_lang}, regardless of the language of the sources you find. "
                        "Be informative — give real content, not just a list of links."
                    ),
                )
                extra = extra_response.output_text.strip()
            except Exception:
                extra = ""

        images = [] if no_answer else self.get_slide_images(metas)
        return {
            "answer": answer,
            "extra": extra,
            "sources": [] if no_answer else sources,
            "images": images,
        }

    def search_online(self, question: str, lang: str = "EN") -> str:
        reply_lang = "Dutch" if lang == "NL" else "English"
        try:
            resp = self.client.responses.create(
                model="gpt-4o-mini",
                tools=[{"type": "web_search_preview"}],
                input=(
                    f"The student asked: \"{question}\". "
                    "Search online and give a concise answer (1-2 sentences max) about this topic. "
                    f"You MUST respond entirely in {reply_lang}. "
                    "Be direct and informative — no links, no lists."
                ),
            )
            return resp.output_text.strip()
        except Exception:
            return ""

    _FOLLOWUP_PHRASES = [
        "explain better", "explain more", "elaborate", "more detail", "more info",
        "tell me more", "can you explain", "what do you mean", "give an example",
        "i don't understand", "i dont understand", "don't get it", "dont get it",
        "leg beter uit", "meer uitleg", "vertel meer", "wat bedoel je", "meer detail",
        "geef een voorbeeld", "kun je uitleggen", "kan je uitleggen", "meer informatie",
        "verduidelijk", "hoe bedoel je", "ik snap het niet", "snap het niet",
        "ik begrijp het niet", "begrijp het niet", "wat betekent dit", "wat betekent dat",
        "niet begrepen", "leg uit", "uitleggen",
    ]

    def _is_followup(self, question: str) -> bool:
        q = question.lower().strip()
        return any(phrase in q for phrase in self._FOLLOWUP_PHRASES)

    def _get_search_query(self, question: str, history: list) -> str:
        if self._is_followup(question) and history:
            for msg in reversed(history):
                if msg.get("role") == "user":
                    return msg["content"]
        return question

    def stream_query(self, question: str, books: list = None, history: list = None, lang: str = "EN", teacher_style: bool = False):
        total = self.collection.count()
        if total == 0:
            yield {"type": "token", "text": "I couldn't find anything — no course materials have been uploaded yet. Please ask your teacher."}
            yield {"type": "done", "sources": []}
            return

        search_query = self._get_search_query(question, history or [])
        response = self.client.embeddings.create(
            input=[self._expand_query(search_query)],
            model="text-embedding-3-small",
        )
        query_embedding = response.data[0].embedding

        where = None
        if books and len(books) > 0:
            if len(books) == 1:
                where = {"filename": {"$eq": books[0]}}
            else:
                where = {"$or": [{"filename": {"$eq": b}} for b in books]}

        n_results = min(2, total)
        results = self.collection.query(query_embeddings=[query_embedding], n_results=n_results, where=where)

        docs = results["documents"][0]
        metas = results["metadatas"][0]

        if not docs:
            yield {"type": "token", "text": "I couldn't find this in your course material, click below to check relevant information online."}
            yield {"type": "done", "sources": []}
            return

        context_parts = []
        sources = []
        for doc, meta in zip(docs, metas):
            fname = meta.get("filename", "Unknown file")
            loc = meta.get("location", "")
            context_parts.append(f"[{fname} — {loc}]\n{doc[:700]}")
            sources.append({"filename": fname, "location": loc, "preview": meta.get("preview", doc[:500])})

        lang_instruction = "Answer in Dutch." if lang == "NL" else "Answer in English."
        fallback_phrase = (
            "Ik kon dit niet vinden in je cursusmateriaal, klik hieronder voor relevante informatie online."
            if lang == "NL" else
            "I couldn't find this in your course material, click below to check relevant information online."
        )
        if teacher_style:
            system_prompt = (
                f"Je bent de leerkracht van dit vak en geeft mondeling uitleg aan een student. {lang_instruction} "
                "Spreek in de EXACTE stijl zoals de leerkracht praat in de cursustranscripties hieronder. "
                "Gebruik dezelfde woordkeuze, zinsstructuur en toon als in de transcripties. "
                "Spreek informeel, direct en enthousiast. Gebruik 'je' en 'jij'. "
                "Geen droge opsommingen — vertel het verhaal achter de stof. "
                "Geen emoji's, geen links, geen externe kennis. "
                f"Als het antwoord er niet in staat, zeg dan exact: \"{fallback_phrase}\""
            )
        else:
            system_prompt = (
                f"You are a study assistant. {lang_instruction} "
                "Answer using ONLY the course material below. No emojis, no links, no outside knowledge. "
                "Use bullet points for lists. Bold key terms with **term**. Be concise. "
                f"If not found, say exactly: \"{fallback_phrase}\""
            )

        messages = [{"role": "system", "content": system_prompt}]
        if history:
            messages.extend(history[-4:])
        messages.append({"role": "user", "content": f"Material:\n{chr(10).join(context_parts)}\n\nQuestion: {question}"})

        stream = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            temperature=0,
            max_tokens=400,
            stream=True,
        )

        full_answer = ""
        for chunk in stream:
            token = chunk.choices[0].delta.content or ""
            if token:
                full_answer += token
                yield {"type": "token", "text": token}

        no_answer = "couldn't find" in full_answer.lower() or "kon dit niet vinden" in full_answer.lower()
        images = [] if no_answer else self.get_slide_images(metas)
        yield {"type": "done", "sources": [] if no_answer else sources, "images": images}

    # ── Search (teacher) ──────────────────────────────────────────────────

    def search_content(self, query: str, course_name: str) -> list:
        total = self.collection.count()
        if total == 0:
            return []

        response = self.client.embeddings.create(
            input=[self._expand_query(query)],
            model="text-embedding-3-small",
        )
        query_embedding = response.data[0].embedding

        where = {"course": {"$eq": course_name}}
        n_results = min(8, total)

        try:
            results = self.collection.query(query_embeddings=[query_embedding], n_results=n_results, where=where)
        except Exception:
            return []

        docs = results["documents"][0]
        metas = results["metadatas"][0]

        return [
            {
                "filename": meta.get("filename", "Unknown file"),
                "location": meta.get("location", ""),
                "preview": doc[:500],
            }
            for doc, meta in zip(docs, metas)
        ]

    # ── Course management ─────────────────────────────────────────────────

    def get_courses(self) -> list:
        if self.collection.count() == 0:
            return []

        results = self.collection.get(include=["metadatas"])
        courses: dict = {}

        for meta in results["metadatas"]:
            course = meta["course"]
            filename = meta["filename"]
            if course not in courses:
                courses[course] = set()
            courses[course].add(filename)

        return [
            {"name": name, "files": sorted(files)}
            for name, files in sorted(courses.items())
        ]

    def delete_course(self, course_name: str):
        results = self.collection.get(where={"course": {"$eq": course_name}}, include=["metadatas"])
        if results["ids"]:
            self.collection.delete(ids=results["ids"])

    def delete_file(self, course_name: str, filename: str):
        results = self.collection.get(
            where={"$and": [{"course": {"$eq": course_name}}, {"filename": {"$eq": filename}}]},
            include=["metadatas"],
        )
        if results["ids"]:
            self.collection.delete(ids=results["ids"])
        if filename.lower().endswith((".pptx", ".ppsx", ".pdf")) and self._sb:
            try:
                rows = (
                    self._sb.table("slide_images")
                    .select("storage_path")
                    .eq("course", course_name)
                    .eq("filename", filename)
                    .execute()
                    .data or []
                )
                for row in rows:
                    try:
                        self._sb.storage.from_("course-files").remove([row["storage_path"]])
                    except Exception:
                        pass
                self._sb.table("slide_images").delete().eq("course", course_name).eq("filename", filename).execute()
            except Exception:
                pass
